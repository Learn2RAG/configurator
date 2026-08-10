"""
generate_test_documents.py

Helper script to generate sample documents for testing all supported file types.
This script creates test documents in the tests/data directory.
"""

import io
import zipfile

from pathlib import Path


def create_pdf() -> bytes:
    """Create a minimal valid PDF file."""
    pdf_content = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj
4 0 obj
<< /Length 100 >>
stream
BT
/F1 12 Tf
50 750 Td
(PDF Document Example) Tj
0 -20 Td
(This is a sample PDF file for testing.) Tj
ET
endstream
endobj
5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj
xref
0 6
0000000000 65535 f 
0000000009 00000 n 
0000000058 00000 n 
0000000115 00000 n 
0000000244 00000 n 
0000000395 00000 n 
trailer
<< /Size 6 /Root 1 0 R >>
startxref
483
%%EOF"""
    return pdf_content


def create_pptx() -> bytes:
    """Create a minimal valid PPTX file using python-pptx."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "PPTX Presentation Example"
    subtitle.text = "This is a sample PPTX file for testing."

    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    return pptx_buffer.getvalue()


def create_docx() -> bytes:
    """Create a minimal valid DOCX file using built-in zip."""
    import zipfile

    docx_buffer = io.BytesIO()
    with zipfile.ZipFile(docx_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        # [Content_Types].xml
        content_types = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
</Types>"""
        zf.writestr("[Content_Types].xml", content_types)

        # _rels/.rels
        rels = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
</Relationships>"""
        zf.writestr("_rels/.rels", rels)

        # word/document.xml
        document = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
<w:body>
<w:p><w:r><w:t>DOCX Document Example</w:t></w:r></w:p>
<w:p><w:r><w:t>This is a sample DOCX file for testing the loaders.</w:t></w:r></w:p>
<w:p><w:r><w:t>It contains multiple paragraphs and formatted text.</w:t></w:r></w:p>
</w:body>
</w:document>"""
        zf.writestr("word/document.xml", document)

        # docProps/core.xml
        core_props = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"/>"""
        zf.writestr("docProps/core.xml", core_props)

    return docx_buffer.getvalue()


def create_xlsx() -> bytes:
    """Create a minimal valid XLSX file in memory using pure zip."""
    xlsx_buffer = io.BytesIO()
    with zipfile.ZipFile(xlsx_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(
            "[Content_Types].xml",
            b'<?xml version="1.0" encoding="UTF-8"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/><Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/></Types>',
        )
        zf.writestr(
            "_rels/.rels",
            b'<?xml version="1.0" encoding="UTF-8"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/></Relationships>',
        )
        zf.writestr(
            "xl/workbook.xml",
            b'<?xml version="1.0" encoding="UTF-8"?><workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheets><sheet name="Sheet1" sheetId="1" r:id="rId1"/></sheets></workbook>',
        )
        zf.writestr(
            "xl/_rels/workbook.xml.rels",
            b'<?xml version="1.0" encoding="UTF-8"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/></Relationships>',
        )
        zf.writestr(
            "xl/worksheets/sheet1.xml",
            b'<?xml version="1.0" encoding="UTF-8"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><sheetData><row r="1"><c r="A1" t="inlineStr"><is><t>Sample Text</t></is></c></row></sheetData></worksheet>',
        )
    return xlsx_buffer.getvalue()


def create_test_documents(data_dir: Path) -> None:
    """Generate all test documents."""
    data_dir.mkdir(parents=True, exist_ok=True)

    # Create PDF
    pdf_path = data_dir / "sample.pdf"
    pdf_path.write_bytes(create_pdf())
    print(f"Created: {pdf_path}")

    # Create XLSX
    xlsx_path = data_dir / "sample.xlsx"
    xlsx_path.write_bytes(create_xlsx())
    print(f"Created: {xlsx_path}")

    # Create DOCX
    docx_path = data_dir / "sample.docx"
    docx_path.write_bytes(create_docx())
    print(f"Created: {docx_path}")

    # Create PPTX
    pptx_path = data_dir / "sample.pptx"
    pptx_path.write_bytes(create_pptx())
    print(f"Created: {pptx_path}")


if __name__ == "__main__":
    data_dir = Path(__file__).parent / "data"
    create_test_documents(data_dir)
    print(f"\nAll test documents created in: {data_dir}")